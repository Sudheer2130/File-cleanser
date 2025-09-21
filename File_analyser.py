import streamlit as st
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import pptx
import docx
import openpyxl
import io
import os
from groq import Groq

# ---------- APP CONFIG ----------
st.set_page_config(
    page_title="VIT AP File Cleanser & Analyzer",
    page_icon="🧹",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------- CUSTOM STYLES ----------
st.markdown(
    """
    <style>
    body {
        background-color: #0D1117;  /* Dark background */
        color: #C9D1D9;             /* Light gray text */
    }
    .stButton>button {
        background-color: #FF6F61;
        color: white;
        border-radius: 10px;
        height: 3em;
    }
    .stDownloadButton>button {
        background-color: #4CAF50;
        color: white;
        border-radius: 10px;
        height: 3em;
    }
    .stTextArea>div>textarea {
        background-color: #161B22;
        color: #C9D1D9;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ---------- VIT AP LOGO ----------
st.image("https://upload.wikimedia.org/wikipedia/en/7/7f/VIT_University_logo.png", width=150)

st.title("📄 VIT AP Automated File Cleansing & Analysis")

# ---------- GROQ API ----------
groq_api_key = os.getenv("GROQ_API_KEY")  # Ensure your key is set in env
client = Groq(api_key=groq_api_key)

# ---------- PRESIDIO ENGINES ----------
analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()

# ---------- FILE UPLOADER ----------
uploaded_file = st.file_uploader(
    "Upload a file",
    type=["txt", "pdf", "png", "jpg", "jpeg", "docx", "xlsx", "pptx"]
)
text = ""

if uploaded_file:
    st.success(f"✅ File uploaded: {uploaded_file.name}")

    # --- TXT FILE ---
    if uploaded_file.type == "text/plain":
        text = uploaded_file.read().decode("utf-8")

    # --- PDF FILE ---
    elif uploaded_file.type == "application/pdf":
        pdf = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page in pdf:
            text += page.get_text()

    # --- DOCX FILE ---
    elif uploaded_file.name.endswith((".docx", ".DOCX")):
        doc = docx.Document(uploaded_file)
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"

    # --- EXCEL FILE ---
    elif uploaded_file.name.endswith((".xlsx", ".XLSX")):
        uploaded_file.seek(0)
        workbook = openpyxl.load_workbook(
            filename=io.BytesIO(uploaded_file.read()), data_only=True
        )
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell) + " "
                text += "\n"

    # --- POWERPOINT FILE ---
    elif uploaded_file.name.endswith((".pptx", ".PPTX")):
        uploaded_file.seek(0)
        presentation = pptx.Presentation(uploaded_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    # --- IMAGE FILE ---
    elif uploaded_file.type.startswith("image"):
        img = Image.open(uploaded_file)
        text = pytesseract.image_to_string(img)

    # ---------- SHOW EXTRACTED TEXT ----------
    if text.strip():
        st.subheader("📜 Extracted Text")
        st.text_area("Raw Content", text, height=200)

        # --- PRESIDIO ANONYMIZATION ---
        results = analyzer.analyze(text=text, language="en")
        anonymized = anonymizer.anonymize(text=text, analyzer_results=results)

        st.subheader("🔒 Anonymized Text")
        st.text_area("Cleansed Output", anonymized.text, height=200)

        # --- DOWNLOAD BUTTON ---
        st.download_button(
            label="⬇️ Download Cleaned Text",
            data=anonymized.text,
            file_name="cleaned_output.txt"
        )

        # --- AI INSIGHTS WITH GROQ ---
        if groq_api_key:
            st.subheader("🤖 AI Insights (Groq)")
            try:
                completion = client.chat.completions.create(
                    model="llama2-70b-4096",  # Updated supported model
                    messages=[
                        {"role": "system", "content": "You are an assistant that summarizes and analyzes text."},
                        {"role": "user", "content": f"Summarize and analyze this text:\n\n{anonymized.text}"}
                    ],
                    max_tokens=300
                )
                ai_response = completion.choices[0].message.content
                st.markdown(
                    f"<div style='background-color:#161B22;padding:15px;border-radius:10px;color:#C9D1D9'>{ai_response}</div>",
                    unsafe_allow_html=True
                )
            except Exception as e:
                st.error(f"Groq API error: {e}\nMake sure you are using a valid model and your API key is correct.")
        else:
            st.info("No GROQ_API_KEY found. Please set it to enable AI insights.")
    else:
        st.warning("⚠️ No text could be extracted from this file.")
