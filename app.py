import streamlit as st
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import pptx
import docx
import openpyxl
import io

# Initialize Presidio engines
analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()

st.title("📄 Automated File Cleansing & Analysis")

# File uploader
uploaded_file = st.file_uploader("Upload a file", type=["txt", "pdf", "png", "jpg", "jpeg", "docx",  "xlsx","pptx"])

text = ""

if uploaded_file:
    st.success(f"✅ File uploaded: {uploaded_file.name}")

    # --- Handle TXT files ---
    if uploaded_file.type == "text/plain":
        text = uploaded_file.read().decode("utf-8")

    # --- Handle PDF files (extract text using PyMuPDF) ---
    elif uploaded_file.type == "application/pdf":
        pdf = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page in pdf:
            text += page.get_text()

    elif  uploaded_file.name.endswith((".docx", ".DOCX")):
        doc = docx.Document(uploaded_file)
        # Extract paragraphs
        for para in doc.paragraphs:
            text += para.text + "\n"
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
    
    # --- Handle Excel Files (.xlsx) ---
    elif uploaded_file.name.endswith((".xlsx", ".XLSX")):
        uploaded_file.seek(0)  # Ensure we're at the start of the file
        workbook = openpyxl.load_workbook(filename=io.BytesIO(uploaded_file.read()), data_only=True)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell) + " "
                text += "\n"

    # --- Handle PowerPoint Files (.pptx) ---
    elif uploaded_file.name.endswith((".pptx", ".PPTX")):
        uploaded_file.seek(0)  # Ensure we're at the start of the file
        presentation = pptx.Presentation(uploaded_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"


    # --- Handle Images (OCR with Pytesseract) ---
    elif uploaded_file.type.startswith("image"):
        img = Image.open(uploaded_file)
        text = pytesseract.image_to_string(img)

    # Show extracted text
    if text.strip():
        st.subheader("📜 Extracted Text")
        st.text_area("Raw Content", text, height=200)

        # --- Run Presidio anonymization ---
        results = analyzer.analyze(text=text, language="en")
        anonymized = anonymizer.anonymize(text=text, analyzer_results=results)

        st.subheader("🔒 Anonymized Text")
        st.text_area("Cleansed Output", anonymized.text, height=200)

        # Download option
        st.download_button(
            label="⬇️ Download Cleaned Text",
            data=anonymized.text,
            file_name="cleaned_output.txt"
        )
    else:
        st.warning("⚠️ No text could be extracted from this file.")
