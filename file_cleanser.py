import streamlit as st
from datetime import datetime
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
import docx
import openpyxl
import io
import os

# LangChain imports
from langchain.llms.base import LLM
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory

# ---------- STREAMLIT CONFIG ----------
st.set_page_config(
    page_title="File Analyzer",
    page_icon="🧹",
    layout="wide"
)

# ---------- STYLING ----------
st.markdown("""
<style>
body {background-color: #0D1117; color: #C9D1D9;}
.stButton>button {background-color: #FF6F61; color: white; border-radius:10px; height:3em;}
.stDownloadButton>button {background-color: #4CAF50; color:white; border-radius:10px; height:3em;}
.stTextArea>div>textarea {background-color:#161B22; color:#C9D1D9;}
</style>
""", unsafe_allow_html=True)

st.title("📄 File Analyzer with AI")

# ---------- SESSION STATE ----------
if "upload_history" not in st.session_state:
    st.session_state.upload_history = []

# ---------- PRESIDIO ----------
analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()

# ---------- GROQ LLM ----------
class GroqLLM(LLM):
    api_key: str
    @property
    def _llm_type(self):
        return "groq"
    def _call(self, prompt: str, stop=None):
        from groq import Groq
        client = Groq(api_key=self.api_key)
        completion = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{"role":"user","content": prompt}],
            max_tokens=500
        )
        return completion.choices[0].message.content

# ---------- FILE UPLOADER ----------
uploaded_files = st.file_uploader(
    "Upload files",
    type=["txt","pdf","docx","xlsx","pptx","png","jpg","jpeg"],
    accept_multiple_files=True
)

# ---------- TEXT EXTRACTION ----------
def extract_text(file):
    text = ""
    if file.type == "text/plain":
        text = file.read().decode("utf-8")
    elif file.type == "application/pdf":
        pdf = fitz.open(stream=file.read(), filetype="pdf")
        for page in pdf: text += page.get_text()
    elif file.name.endswith((".docx", ".DOCX")):
        doc = docx.Document(file)
        for para in doc.paragraphs: text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells: text += cell.text + " "
                text += "\n"
    elif file.name.endswith((".xlsx", ".XLSX")):
        file.seek(0)
        workbook = openpyxl.load_workbook(io.BytesIO(file.read()), data_only=True)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell: text += str(cell) + " "
                text += "\n"
    elif file.name.endswith((".pptx", ".PPTX")):
        file.seek(0)
        pres = Presentation(file)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text"): text += shape.text + "\n"
    elif file.type.startswith("image"):
        img = Image.open(file)
        text = pytesseract.image_to_string(img)
    return text

# ---------- PROCESS FILES ----------
if uploaded_files:
    groq_api_key = os.getenv("GROQ_API_KEY")  # Or hardcode for testing
    llm = GroqLLM(api_key=groq_api_key)
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    for file in uploaded_files:
        original_text = extract_text(file)

        st.session_state.upload_history.append({
            "filename": file.name,
            "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        })

        # ---------- ANONYMIZE ----------
        results = analyzer.analyze(text=original_text, language="en")
        anonymized_text = anonymizer.anonymize(text=original_text, analyzer_results=results).text

        # ---------- AI SUMMARIZATION ----------
        prompt = PromptTemplate(
            input_variables=["text"],
            template="You are an AI assistant. Summarize the following text concisely:\n\n{text}"
        )
        chain = LLMChain(llm=llm, prompt=prompt, memory=memory)
        try:
            ai_summary = chain.run(text=anonymized_text)
        except:
            ai_summary = "AI summary unavailable."

        st.subheader(f"🤖 AI Summary for {file.name}")
        st.text_area("AI Summary", ai_summary, height=200)

        # ---------- EXPORT CLEAN FILE ----------
        if file.name.endswith(".txt"):
            st.download_button(
                label=f"⬇️ Download Cleaned TXT ({file.name})",
                data=anonymized_text,
                file_name=f"cleaned_{file.name}"
            )

        elif file.name.endswith((".docx", ".DOCX")):
            doc = docx.Document()
            for line in anonymized_text.split("\n"):
                doc.add_paragraph(line)
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            st.download_button(
                label=f"⬇️ Download Cleaned DOCX ({file.name})",
                data=buffer,
                file_name=f"cleaned_{file.name}"
            )

        elif file.name.endswith((".xlsx", ".XLSX")):
            # Preserve structure of original workbook
            file.seek(0)
            original_wb = openpyxl.load_workbook(io.BytesIO(file.read()), data_only=True)
            clean_wb = openpyxl.Workbook()
            clean_wb.remove(clean_wb.active)

            for sheet in original_wb.worksheets:
                new_sheet = clean_wb.create_sheet(title=sheet.title)
                for r, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    row_text = " ".join([str(cell) if cell else "" for cell in row])
                    results = analyzer.analyze(text=row_text, language="en")
                    anonymized_row_text = anonymizer.anonymize(text=row_text, analyzer_results=results).text
                    # Keep proper columns
                    for c, cell_val in enumerate(anonymized_row_text.split(), start=1):
                        new_sheet.cell(row=r, column=c, value=cell_val)

            buffer = io.BytesIO()
            clean_wb.save(buffer)
            buffer.seek(0)
            st.download_button(
                label=f"⬇️ Download Cleaned XLSX ({file.name})",
                data=buffer,
                file_name=f"cleaned_{file.name}",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

        elif file.name.endswith((".pptx", ".PPTX")):
            pres = Presentation()
            for line in anonymized_text.split("\n"):
                slide = pres.slides.add_slide(pres.slide_layouts[5])
                slide.shapes.title.text = line[:100]
            buffer = io.BytesIO()
            pres.save(buffer)
            buffer.seek(0)
            st.download_button(
                label=f"⬇️ Download Cleaned PPTX ({file.name})",
                data=buffer,
                file_name=f"cleaned_{file.name}",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        elif file.type.startswith("image") or file.name.endswith(".pdf"):
            st.download_button(
                label=f"⬇️ Download Cleaned TXT ({file.name})",
                data=anonymized_text,
                file_name=f"cleaned_{file.name}.txt"
            )

# ---------- DISPLAY UPLOAD HISTORY ----------
st.subheader("📜 Upload History")
for i, entry in enumerate(st.session_state.upload_history[::-1], start=1):
    st.write(f"**{i}. {entry['filename']}** — uploaded at *{entry['timestamp']}*")
